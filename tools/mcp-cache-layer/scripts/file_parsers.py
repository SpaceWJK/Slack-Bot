"""
file_parsers.py — 원본 파일 직접 파싱 고도화 모듈

GDI MCP 청크 텍스트 대신 원본 파일(PPTX, XLSX, TSV, PNG)에서 직접 파싱하여
사람이 보는 형태와 최대한 동일한 마크다운 텍스트를 생성한다.

사용법:
    from file_parsers import parse_file
    result = parse_file("path/to/file.pptx")
    # result = {"body_text": "## Slide 1\n...", "metadata": {...}, "images": [...]}

지원 형식:
    - PPTX: 슬라이드별 텍스트 + 테이블 + 이미지 참조
    - XLSX: 시트별 마크다운 테이블 (병합 셀 fill-down, 빈 행 스킵)
    - TSV:  원본 탭 구분 직접 파싱 → 마크다운 테이블
    - PNG:  메타데이터 + (선택) Vision API 분석

의존성:
    pip install python-pptx openpyxl Pillow pytesseract
    Tesseract OCR 바이너리 별도 설치 필요 (https://github.com/UB-Mannheim/tesseract/wiki)
    미설치 시 PPTX 이미지 OCR 건너뜀 (기존 동작 유지)
"""

import csv
import io
import logging
import os
import re
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)

# ── 상수 ──────────────────────────────────────────────────────────────────────
MAX_TABLE_ROWS = 20000      # 마크다운 테이블 최대 행 수
MAX_BODY_CHARS = 500_000    # body_text 최대 글자 수
MAX_COL_WIDTH = 80          # 셀 값 최대 표시 길이 (넘으면 잘라냄)
MAX_PRE_HEADER_CELLS_PER_ROW = 50  # task-127 S4-3: pre_header 행당 셀 최대 출력 수 (폭증 방지)
# task-127 v3 시정 2: multi-line 셀 끝에 issue ID 패턴이 있으면 절단 시 보존
# (#162733/#164993/#167702 같은 형식. 4자리 이상 숫자 — false positive 회피)
_TRAILING_ID_RE = re.compile(r"#\d{4,}\b")
IMAGE_DIR_NAME = "_images"  # 이미지 추출 저장 폴더명
_tesseract_available: "bool | None" = None  # None=미확인, True=가용, False=불가


# ═══════════════════════════════════════════════════════════════════════════════
# 통합 진입점
# ═══════════════════════════════════════════════════════════════════════════════

def parse_file(file_path: str, *, extract_images: bool = True,
               image_base_dir: Optional[str] = None) -> dict:
    """원본 파일을 파싱하여 구조화된 결과를 반환한다.

    Args:
        file_path: 원본 파일 경로
        extract_images: PPTX 이미지를 디스크에 추출할지 여부
        image_base_dir: 이미지 저장 기본 경로 (None이면 파일 옆 _images/)

    Returns:
        {
            "body_text": str,      # 마크다운 텍스트
            "metadata": dict,      # 파일 메타데이터
            "images": list[dict],  # 추출된 이미지 정보 (PPTX만)
            "source_type": str,    # generic_pptx / generic_xlsx / generic_tsv / image
        }
    """
    ext = Path(file_path).suffix.lower()
    if ext == ".pptx":
        return parse_pptx(file_path, extract_images=extract_images,
                          image_base_dir=image_base_dir)
    elif ext in (".xlsx", ".xlsm"):
        return parse_xlsx(file_path, extract_images=extract_images,
                          image_base_dir=image_base_dir)
    elif ext == ".pdf":
        return parse_pdf(file_path)
    elif ext == ".tsv":
        return parse_tsv(file_path)
    elif ext in (".docx",):
        return parse_docx(file_path)
    elif ext in (".html", ".htm"):
        return parse_html(file_path, extract_images=extract_images,
                         image_base_dir=image_base_dir)
    elif ext in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"):
        return parse_image(file_path)
    else:
        return {
            "body_text": f"[지원하지 않는 형식: {ext}]",
            "metadata": {"file_path": file_path, "ext": ext},
            "images": [],
            "source_type": "unknown",
        }


# ═══════════════════════════════════════════════════════════════════════════════
# PPTX 파서
# ═══════════════════════════════════════════════════════════════════════════════

def parse_pptx(file_path: str, *, extract_images: bool = True,
               image_base_dir: Optional[str] = None) -> dict:
    """PPTX 파일을 파싱하여 슬라이드별 마크다운을 생성한다.

    각 슬라이드를 ## Slide N 헤더로 구분하고:
    - 텍스트: 그대로 출력
    - 테이블: 마크다운 테이블로 변환
    - 이미지: [이미지: 파일명, 크기] 참조 + 디스크 저장
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(file_path)
    file_name = Path(file_path).stem
    parts = []
    images = []
    total_images = 0

    # 이미지 저장 디렉토리
    if extract_images:
        if image_base_dir:
            img_dir = Path(image_base_dir) / IMAGE_DIR_NAME / _safe_filename(file_name)
        else:
            img_dir = Path(file_path).parent / IMAGE_DIR_NAME / _safe_filename(file_name)
    else:
        img_dir = None

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_parts = []
        slide_parts.append(f"## Slide {slide_idx}")

        # 슬라이드 내 shape을 위치 순서(top→left)로 정렬
        shapes = sorted(slide.shapes, key=lambda s: (s.top or 0, s.left or 0))

        for shape in shapes:
            # ── 텍스트 프레임 ──
            if shape.has_text_frame:
                text = _extract_text_frame(shape.text_frame)
                if text:
                    slide_parts.append(text)

            # ── 테이블 ──
            if shape.has_table:
                table_md = _extract_table(shape.table)
                if table_md:
                    slide_parts.append(table_md)

            # ── 이미지 ──
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                total_images += 1
                img_info = _extract_image(shape, slide_idx, total_images,
                                          img_dir, extract_images)
                if img_info:
                    images.append(img_info)
                    size_kb = img_info["size_bytes"] // 1024
                    ref = f'[이미지: {img_info["filename"]}, {size_kb}KB]'
                    slide_parts.append(ref)
                    if img_info.get("saved_path"):
                        ocr_text = _ocr_image(img_info["saved_path"])
                        if ocr_text:
                            slide_parts.append(ocr_text)

            # ── 그룹 셰이프 (내부 텍스트 추출) ──
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_text = _extract_group_text(shape)
                if group_text:
                    slide_parts.append(group_text)

        # ── 발표자 노트 (speaker notes) ──
        # 슬라이드 본문에 나타나지 않지만 기획/검토/QA 맥락이 다수 포함됨
        if slide.has_notes_slide:
            try:
                notes_frame = slide.notes_slide.notes_text_frame
                notes_text = _extract_text_frame(notes_frame)
                if notes_text and notes_text.strip():
                    slide_parts.append(f"> **[발표자 노트]** {notes_text}")
            except Exception as e:
                logger.debug("슬라이드 %d 노트 추출 실패: %s", slide_idx, e)

        # 빈 슬라이드 건너뛰기
        content = "\n".join(slide_parts[1:])  # 헤더 제외 내용
        if not content.strip():
            continue

        parts.append("\n".join(slide_parts))

    body_text = "\n\n".join(parts)
    if len(body_text) > MAX_BODY_CHARS:
        body_text = body_text[:MAX_BODY_CHARS] + "\n\n_(본문 잘림)_"

    return {
        "body_text": body_text,
        "metadata": {
            "file_path": file_path,
            "file_name": Path(file_path).name,
            "total_slides": len(prs.slides),
            "total_images": total_images,
            "slide_width": prs.slide_width,
            "slide_height": prs.slide_height,
        },
        "images": images,
        "source_type": "generic_pptx",
    }


def _extract_text_frame(text_frame) -> str:
    """텍스트 프레임에서 문단별 텍스트를 추출한다."""
    lines = []
    for para in text_frame.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        # 글머리 기호 처리
        if para.level and para.level > 0:
            indent = "  " * para.level
            text = f"{indent}- {text}"
        lines.append(text)
    return "\n".join(lines)


def _extract_table(table) -> str:
    """PPTX 테이블을 마크다운 테이블로 변환한다.

    task-127 v3 시정 3b: cell text 절단 시 끝 ID 패턴 보존 (xlsx _sanitize_cells 정합).
    """
    rows = []
    for r_idx in range(len(table.rows)):
        cells = []
        for c_idx in range(len(table.columns)):
            cell_text = table.cell(r_idx, c_idx).text.strip()
            cell_text = cell_text.replace("|", "\\|").replace("\n", " ")
            if len(cell_text) > MAX_COL_WIDTH:
                # task-127 v3 시정 3b: 끝부분 ID 패턴 보존
                tail_ids = _TRAILING_ID_RE.findall(cell_text)
                if tail_ids:
                    tail_str = " " + " ".join(dict.fromkeys(tail_ids))
                    head_w = max(0, MAX_COL_WIDTH - len(tail_str))
                    cell_text = cell_text[:head_w] + "…" + tail_str
                else:
                    cell_text = cell_text[:MAX_COL_WIDTH] + "…"
            cells.append(cell_text)
        rows.append(cells)

    if not rows:
        return ""

    result = []
    # 첫 행 = 헤더
    result.append("| " + " | ".join(rows[0]) + " |")
    result.append("|" + "|".join("---" for _ in rows[0]) + "|")
    for row in rows[1:]:
        result.append("| " + " | ".join(row) + " |")

    return "\n".join(result)


_TESSERACT_FALLBACK_PATHS = (
    r"C:\Program Files\Tesseract-OCR\tesseract.exe",
    r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
)
_TESSDATA_FALLBACK_PATHS = (
    os.path.expandvars(r"%USERPROFILE%\tessdata"),
    r"C:\Program Files\Tesseract-OCR\tessdata",
    r"C:\Program Files (x86)\Tesseract-OCR\tessdata",
)


def _autodetect_tesseract() -> None:
    """Tesseract 바이너리 + tessdata 디렉터리를 자동 탐지하여 pytesseract/환경변수 설정.

    Windows winget 설치 직후에는 현 프로세스 PATH에 반영이 안 될 수 있으므로,
    표준 설치 경로를 순차 탐색한다. 한국어 언어팩이 `%USERPROFILE%\\tessdata`에
    있는 경우도 지원한다 (관리자 권한 없이 설치 가능한 경로).
    """
    import pytesseract as _pt

    # tesseract 바이너리
    if not getattr(_pt.pytesseract, "_autodetected", False):
        for p in _TESSERACT_FALLBACK_PATHS:
            if os.path.isfile(p):
                _pt.pytesseract.tesseract_cmd = p
                break
        _pt.pytesseract._autodetected = True

    # tessdata 디렉터리 (kor.traineddata 우선 확인)
    if not os.environ.get("TESSDATA_PREFIX"):
        for p in _TESSDATA_FALLBACK_PATHS:
            if os.path.isfile(os.path.join(p, "kor.traineddata")):
                os.environ["TESSDATA_PREFIX"] = p
                break


def _ocr_image(image_path: str, lang: str = "kor+eng") -> str:
    """Tesseract OCR로 이미지에서 텍스트를 추출한다.

    Tesseract 미설치 시 경고 1회 로그 후 빈 문자열 반환 (graceful degradation).
    Windows 표준 설치 경로를 자동 탐지하므로 시스템 PATH 갱신 없이도 동작.
    """
    global _tesseract_available
    if _tesseract_available is False:
        return ""
    try:
        import pytesseract  # 지연 import — 모듈 레벨 ImportError 회피
        from PIL import Image as _PILImage
        if _tesseract_available is None:
            _autodetect_tesseract()
            pytesseract.get_tesseract_version()  # 바이너리 존재 1회 확인
            _tesseract_available = True
        text = pytesseract.image_to_string(_PILImage.open(image_path), lang=lang)
        return text.strip()
    except Exception as e:
        if _tesseract_available is None:
            _tesseract_available = False
            logger.warning("[file_parsers] Tesseract 미설치 — PPTX OCR 건너뜀: %s", e)
        else:
            logger.debug("[file_parsers] OCR 실패 (%s): %s", Path(image_path).name, e)
        return ""


def _extract_image(shape, slide_idx: int, img_seq: int,
                   img_dir: Optional[Path], save: bool) -> Optional[dict]:
    """이미지 셰이프에서 정보를 추출하고 선택적으로 디스크에 저장한다."""
    try:
        img = shape.image
        ext = img.ext or "png"
        blob = img.blob
        content_type = img.content_type
        filename = f"slide{slide_idx:02d}_img{img_seq:03d}.{ext}"

        saved_path = None
        if save and img_dir:
            img_dir.mkdir(parents=True, exist_ok=True)
            save_path = img_dir / filename
            save_path.write_bytes(blob)
            saved_path = str(save_path)

        return {
            "filename": filename,
            "content_type": content_type,
            "size_bytes": len(blob),
            "slide": slide_idx,
            "saved_path": saved_path,
            "width": shape.width,
            "height": shape.height,
        }
    except Exception as e:
        logger.warning("이미지 추출 실패 (slide %d): %s", slide_idx, e)
        return None


def _extract_group_text(group_shape) -> str:
    """그룹 셰이프 내부의 텍스트 + 표 + 중첩 그룹을 재귀적으로 추출한다.

    task-127 v3 시정 3a: GROUP 내 TABLE 처리 추가 (송윤선 pptx Slide 6/7/9 한글 셀 누락 시정).
    """
    texts = []
    for shape in group_shape.shapes:
        # text frame
        if getattr(shape, "has_text_frame", False):
            t = _extract_text_frame(shape.text_frame)
            if t:
                texts.append(t)
        # task-127 v3 시정 3a: GROUP 내 TABLE 처리 추가
        if getattr(shape, "has_table", False):
            md = _extract_table(shape.table)
            if md:
                texts.append(md)
        # 중첩 그룹 (재귀 진입 — has_text_frame/has_table만 있는 일반 shape는 .shapes 미보유)
        # MSO_SHAPE_TYPE.GROUP 직접 비교 대신 .shapes attribute 존재 + iterable 검증
        inner_shapes = getattr(shape, "shapes", None)
        if inner_shapes is not None and shape is not group_shape:
            try:
                # iterable 검증 (text_frame/table 같은 일반 shape의 .shapes는 mock 시 MagicMock 반환)
                iter(inner_shapes)
                # group이면 재귀, 일반 shape의 spurious .shapes면 빈 결과
                t = _extract_group_text(shape)
                if t:
                    texts.append(t)
            except TypeError:
                pass
    return "\n".join(texts)


# ═══════════════════════════════════════════════════════════════════════════════
# XLSX 파서
# ═══════════════════════════════════════════════════════════════════════════════

def parse_xlsx(file_path: str, *, extract_images: bool = True,
               image_base_dir: Optional[str] = None) -> dict:
    """XLSX 파일을 시트별 마크다운 테이블로 변환한다.

    핵심 처리:
    - 수직 병합: fill-down (카테고리 열)
    - 수평 병합: 중복 열 제거
    - Summary/Report 시트: key-value 형태로 파싱
    - 헤더 감지: 'ID', '번호', '분류' 등 키워드 기반 + 고유값 비율 판별
    - 빈 스페이서 열/행 자동 제거
    - R-1 (task-125): xlsx 이미지 proxy 마커 추출 (extract_images=True 시 디스크 저장)
    - R-2 (task-125): 수식 원문 보존 (data_only=False 보조 로드로 formula_map 추출)

    Args:
        file_path: xlsx 파일 경로
        extract_images: True이면 이미지를 디스크에 저장 (기본값). False이면 마커만 본문에 포함.
        image_base_dir: 이미지 저장 기본 경로 (None이면 파일 옆 _images/)
    """
    import openpyxl

    wb = openpyxl.load_workbook(file_path, data_only=True)
    file_name = Path(file_path).name
    parts = []
    images: list[dict] = []       # R-1 (task-125): 이미지 info 누적
    total_images = 0               # R-1 (task-125): 이미지 수 집계
    total_rows = 0

    # R-1 (task-125): 이미지 저장 디렉터리 결정
    if extract_images:
        if image_base_dir:
            img_dir = Path(image_base_dir) / IMAGE_DIR_NAME / _safe_filename(Path(file_path).stem)
        else:
            img_dir = Path(file_path).parent / IMAGE_DIR_NAME / _safe_filename(Path(file_path).stem)
    else:
        img_dir = None

    # R-2 (task-125): formula map 1회 추출 (read_only 보조 로드 — 옵션 D)
    formula_map_all = _load_xlsx_formula_map(file_path)

    for sname in wb.sheetnames:
        ws = wb[sname]

        # ── R-1 (task-125): 이미지 처리 — C-3 시정: summary/non-summary 모두 처리 ──
        # continue 이전에 실행하여 Summary 시트 이미지 누락 방지
        marker_lines, img_count = _extract_xlsx_image_section(
            ws, sname, img_dir, extract_images, images, total_images
        )
        total_images += img_count

        # Summary/Report 시트 → 워크시트에서 직접 key-value 파싱
        if _is_summary_sheet_ws(sname, ws):
            # task-127 S4-4 MI-1: formula_map 전달 → summary 시트도 수식 augment
            summary_md = _parse_summary_sheet_ws(sname, ws, formula_map_all.get(sname))
            if summary_md:
                parts.append(summary_md)
            # Summary 시트의 이미지 marker도 본문에 추가
            if marker_lines:
                parts.extend(marker_lines)
            continue

        # 수평 병합 범위 수집 (중복 열 제거용)
        h_merges = _collect_horizontal_merges(ws)

        # R-2 (task-125): formula_map 시트별 추출
        sheet_formula_map = formula_map_all.get(sname)

        # 원본 셀 값 읽기 (병합 해제 안 함 → 수평 중복 방지)
        raw_rows, use_cols = _read_sheet_raw(ws, h_merges, sheet_formula_map)

        if not raw_rows:
            # non-summary 시트에도 이미지 marker 추가 (데이터 없어도)
            if marker_lines:
                parts.extend(marker_lines)
            continue

        # 빈 열 제거 (전체가 빈 열)
        col_count = len(raw_rows[0])
        non_empty_cols = []
        for c in range(col_count):
            if any(r[c] for r in raw_rows):
                non_empty_cols.append(c)

        if not non_empty_cols:
            if marker_lines:
                parts.extend(marker_lines)
            continue

        # 컬럼 영역 분리 (원본 워크시트 열 기준 3열 이상 갭이면 별도 테이블)
        regions = _split_column_regions(non_empty_cols, use_cols, raw_rows)
        sheet_header_added = False

        for region_idx, region_cols in enumerate(regions):
            filtered_rows = [[row[c] for c in region_cols] for row in raw_rows]

            # 모든 행이 비어있으면 건너뛰기
            if not any(any(v for v in r) for r in filtered_rows):
                continue

            # 영역 내 빈 열 제거 (대부분 비어있는 열)
            region_col_count = len(region_cols)
            keep_in_region = []
            for rc in range(region_col_count):
                if any(row[rc] for row in filtered_rows):
                    keep_in_region.append(rc)
            if keep_in_region and len(keep_in_region) < region_col_count:
                filtered_rows = [[row[rc] for rc in keep_in_region]
                                 for row in filtered_rows]

            # 테이블 헤더 행 감지
            # task-127 S4-2 NC-1 시정: region_cols<=2인 단일/이중 컬럼에서 header_idx=None이면
            # 첫 행을 헤더로 쓰면 데이터 1행 손실 → 더미 헤더 채택
            header_idx = _detect_header_row(filtered_rows)
            if header_idx is None:
                if len(region_cols) <= 2:
                    # NC-1: 더미 헤더 "(값)" × N — 데이터 전체 보존
                    pre_header = []
                    headers = ["(값)"] * len(region_cols)
                    data_rows = filtered_rows
                else:
                    header_idx = 0
                    pre_header = filtered_rows[:header_idx]
                    headers = filtered_rows[header_idx]
                    data_rows = filtered_rows[header_idx + 1:]
            else:
                # 헤더 전 행은 메타데이터로 추출
                pre_header = filtered_rows[:header_idx]
                headers = filtered_rows[header_idx]
                data_rows = filtered_rows[header_idx + 1:]

            # 중복 헤더 제거 (수평 병합으로 같은 텍스트가 여러 열에)
            headers, keep_cols = _deduplicate_headers(headers)
            if keep_cols:
                data_rows = [[r[c] for c in keep_cols] for r in data_rows]

            # 빈 행 건너뛰기
            data_rows = [r for r in data_rows if any(v for v in r)]

            if not data_rows and not any(v for v in headers):
                continue

            # 카테고리 fill-down (빈 셀을 위 행 값으로 채움)
            if data_rows:
                data_rows = _fill_down_categories(headers, data_rows)

            # 시트 헤더 출력 (첫 영역에서만)
            total_rows += len(data_rows)
            if not sheet_header_added:
                parts.append(f"## Sheet: {sname}")
                sheet_header_added = True
            elif region_idx > 0:
                parts.append("")  # 영역 구분

            # 헤더 전 요약 정보 (있으면)
            meta_text = _format_pre_header(pre_header)
            if meta_text:
                parts.append(meta_text)
                parts.append("")

            # 마크다운 테이블
            if data_rows or any(v for v in headers):
                sanitized_headers = _sanitize_cells(headers)
                parts.append("| " + " | ".join(sanitized_headers) + " |")
                parts.append("|" + "|".join("---" for _ in sanitized_headers) + "|")

                for i, row in enumerate(data_rows):
                    if i >= MAX_TABLE_ROWS:
                        parts.append(
                            f"\n_(… 외 {len(data_rows) - MAX_TABLE_ROWS}행 생략)_"
                        )
                        break
                    # 열 수 맞추기
                    while len(row) < len(headers):
                        row.append("")
                    parts.append(
                        "| " + " | ".join(_sanitize_cells(row[:len(headers)])) + " |"
                    )

        # R-1 (task-125): non-summary 시트 끝에 이미지 marker 추가
        if marker_lines:
            parts.extend(marker_lines)

        parts.append("")

    body_text = "\n".join(parts)
    if len(body_text) > MAX_BODY_CHARS:
        body_text = body_text[:MAX_BODY_CHARS] + "\n\n_(본문 잘림)_"

    return {
        "body_text": body_text,
        "metadata": {
            "file_path": file_path,
            "file_name": file_name,
            "sheet_names": wb.sheetnames,
            "total_sheets": len(wb.sheetnames),
            "total_data_rows": total_rows,
            "total_images": total_images,  # R-1 (task-125): 신규
        },
        "images": images,              # R-1 (task-125): [] → list[info]
        "source_type": "generic_xlsx",
    }


def _collect_horizontal_merges(ws) -> dict:
    """수평 병합 범위를 {(row, col): (min_row, min_col)} 매핑으로 수집.

    수평 병합: 같은 행 내 여러 열이 하나의 값을 공유.
    수직 병합은 별도로 fill-down 처리하므로 여기서는 수평만.
    """
    h_merge_cells = {}  # (row, col) → (min_row, min_col)  (첫 셀 제외)
    for merged_range in ws.merged_cells.ranges:
        if merged_range.min_row == merged_range.max_row:
            # 수평 병합 — 첫 열 외 나머지를 "병합됨"으로 표시
            for col in range(merged_range.min_col + 1, merged_range.max_col + 1):
                h_merge_cells[(merged_range.min_row, col)] = (
                    merged_range.min_row, merged_range.min_col
                )
    return h_merge_cells


def _build_vertical_merge_map(ws) -> dict:
    """수직 병합 범위를 (row, col) → (min_row, min_col) 매핑으로 변환.

    수직 병합: 여러 행이 같은 열에서 하나의 값을 공유 (카테고리 fill-down).
    """
    v_merge_map = {}
    for merged_range in ws.merged_cells.ranges:
        if merged_range.max_row > merged_range.min_row:
            min_row = merged_range.min_row
            min_col = merged_range.min_col
            for row in range(merged_range.min_row + 1, merged_range.max_row + 1):
                for col in range(merged_range.min_col, merged_range.max_col + 1):
                    v_merge_map[(row, col)] = (min_row, min_col)
    return v_merge_map


def _read_sheet_raw(ws, h_merges: dict, formula_map: Optional[dict] = None) -> tuple[list[list[str]], list[int]]:
    """시트의 원본 셀 값을 읽되, 수평 병합 셀은 건너뛴다.

    수직 병합은 첫 셀 값을 유지 (fill-down은 나중에 처리).

    Returns:
        (raw_rows, use_cols) — use_cols는 원본 워크시트 열 번호 리스트
    """
    v_merge_map = _build_vertical_merge_map(ws)

    # 수평 병합으로 인한 중복 열 번호 수집
    h_merged_cols_per_row = {}
    for (r, c) in h_merges:
        h_merged_cols_per_row.setdefault(r, set()).add(c)

    # 전체 열 중 수평 병합으로 항상 숨겨야 할 열 판별
    # (헤더 행에서 수평 병합된 열은 전체 데이터에서도 제거)
    all_skip_cols = set()
    for (r, c) in h_merges:
        all_skip_cols.add(c)

    # 실제로 사용할 열 목록 결정
    all_cols = list(range(1, ws.max_column + 1))
    # 수평 병합 열이 전체의 대부분이면 제거, 아니면 유지
    # → 보수적: 2개 이상의 행에서 병합된 열만 제거
    col_merge_count = {}
    for (r, c) in h_merges:
        col_merge_count[c] = col_merge_count.get(c, 0) + 1

    # task-127 S4-1 옵션 #2: 적응형 임계 + 데이터 셀 보존 가드 (설계 §4.1 v2)
    # total_data_rows: ws.max_row 사용 (1차 증거 — U1b/U2b 단위 테스트 기준)
    # 설계 MAJOR 2의 merged_rows 근사는 mock ws에서 오작동 → ws.max_row 직접 사용
    total_data_rows = max(ws.max_row, 1)

    # col별 데이터 셀 수 사전 계산 (병합 후속 셀 제외)
    h_merges_set = set(h_merges.keys())
    col_data_cells = {}
    for c in range(1, ws.max_column + 1):
        cnt_data = 0
        for r in range(1, ws.max_row + 1):
            if (r, c) in h_merges_set:
                continue  # 병합 후속 셀 — 첫 셀에서만 카운트
            v = ws.cell(row=r, column=c).value
            if v is not None and str(v).strip() and str(v).strip() != "　":
                cnt_data += 1
        col_data_cells[c] = cnt_data

    # 적응형 임계 + 데이터 셀 보존 가드 (옵션 #2)
    adaptive_skip = set()
    for c, cnt in col_merge_count.items():
        if cnt < 2:
            continue
        # 가드 1: 적응형 50% — 구조적 매트릭스 보존 (cnt/max_r >= 0.5이면 skip 제외)
        if cnt / total_data_rows >= 0.5:
            continue
        # 가드 2: 데이터 셀 보존 — col에 데이터 1건이라도 있으면 skip 제외
        if col_data_cells.get(c, 0) >= 1:
            continue
        adaptive_skip.add(c)
    skip_cols = adaptive_skip

    use_cols = [c for c in all_cols if c not in skip_cols]

    raw_rows = []
    for r_idx in range(1, ws.max_row + 1):
        row_vals = []
        for c_idx in use_cols:
            # 수직 병합 → 첫 셀 값 참조
            actual_r, actual_c = v_merge_map.get((r_idx, c_idx), (r_idx, c_idx))
            cell = ws.cell(row=actual_r, column=actual_c)
            val = cell.value
            if val is None:
                val = ""
            else:
                val = str(val).strip()
                if val == "\u3000":
                    val = ""
            # R-2 (task-125): \uc218\uc2dd augment \u2014 formula_map \uc870\ud68c
            if formula_map is not None:
                formula = formula_map.get((actual_r, actual_c))
                val = _format_cell_value_with_formula(val, formula)
            row_vals.append(val)
        raw_rows.append(row_vals)

    return raw_rows, use_cols


def _is_summary_sheet_ws(name: str, ws) -> bool:
    """Summary/Report 형태의 시트인지 판별한다 (워크시트 직접 검사).

    판단 기준:
    - 시트 이름에 'summary', 'report', '요약' 포함
    - 또는: 행 수 < 35이고, 병합 셀이 많고 데이터가 sparse
    """
    name_lower = name.lower()
    if any(kw in name_lower for kw in ('summary', 'report', '요약')):
        return True
    if ws.max_row and ws.max_row <= 35:
        merge_count = len(list(ws.merged_cells.ranges))
        if merge_count > ws.max_row * 0.5:
            return True
    return False


def _parse_summary_sheet_ws(name: str, ws, formula_map: Optional[dict] = None) -> str:
    """Summary 시트를 워크시트에서 직접 key-value 마크다운으로 변환한다.

    QA Report Summary 시트 구조:
    - 섹션 제목 (QA Report, Test Imformation, 테스트 기종, Final Result)
    - key: value 쌍 (제목, 담당자, Client 버전 등)
    - 테이블형 데이터 (OS/Device/OS ver)
    """
    parts = [f"## Sheet: {name}"]

    # 수직 병합 맵 생성 (카테고리 fill을 위해)
    v_merge_map = _build_vertical_merge_map(ws)

    # 테이블 영역 감지 (OS | Device | OS ver 같은 소테이블)
    table_start = None
    table_headers = []

    for r in range(1, ws.max_row + 1):
        cells = []
        for c in range(1, ws.max_column + 1):
            # 수직 병합 → 첫 셀 값 참조
            actual_r, actual_c = v_merge_map.get((r, c), (r, c))
            val = ws.cell(row=actual_r, column=actual_c).value
            if val is not None:
                val = str(val).strip().replace("\n", " ")
                if val and val != "\u3000":
                    # task-127 S4-4: formula_map augment \uc801\uc6a9 (summary \uc2dc\ud2b8\ub3c4 \uc218\uc2dd \ubcf4\uc874)
                    if formula_map is not None:
                        formula = formula_map.get((actual_r, actual_c))
                        val = _format_cell_value_with_formula(val, formula)
                    cells.append((c, val))

        if not cells:
            continue

        # 테이블 모드 진행 중
        if table_start is not None:
            if cells:
                # 1개 값만 있을 때: 수직 병합 fill인지 실제 값인지 구분
                if len(cells) == 1:
                    single_col = cells[0][0]
                    single_val = cells[0][1]
                    # 수직 병합에서 채워진 값이면 → 아티팩트, 건너뛰기
                    actual_r, _ = v_merge_map.get((r, single_col), (r, single_col))
                    if actual_r != r:
                        continue
                    # 실제 셀 값이면 → 섹션 헤더, 테이블 종료
                    table_start = None
                    parts.append("")
                    parts.append(f"\n**{single_val}**")
                    continue

                row_vals = {}
                for col, val in cells:
                    best_h = None
                    best_dist = 999
                    for h_col, h_name in table_headers:
                        dist = abs(col - h_col)
                        if dist < best_dist:
                            best_dist = dist
                            best_h = h_name
                    if best_h and best_dist <= 2:
                        row_vals[best_h] = val

                if row_vals and any(v.strip() for v in row_vals.values()):
                    line = " | ".join(
                        f"{row_vals.get(h, '')}" for _, h in table_headers
                    )
                    parts.append(f"| {line} |")
                else:
                    table_start = None
                    parts.append("")
            continue

        # 1개 값 → 섹션 제목 또는 키워드
        if len(cells) == 1:
            parts.append(f"\n**{cells[0][1]}**")
            continue

        # 2개 값 → key: value
        if len(cells) == 2:
            parts.append(f"- **{cells[0][1]}**: {cells[1][1]}")
            continue

        # 3개 이상 → 소테이블 헤더 or 여러 key:value 쌍
        # 소테이블 감지: 모든 값이 짧은 텍스트(레이블)인 경우
        vals = [v for _, v in cells]
        all_labels = all(len(v) <= 15 and not re.match(r'^[\d,.]+$', v) for v in vals)

        if all_labels and len(cells) >= 3:
            # 소테이블 시작
            table_start = r
            table_headers = cells
            h_line = " | ".join(h for _, h in table_headers)
            sep = " | ".join("---" for _ in table_headers)
            parts.append(f"\n| {h_line} |")
            parts.append(f"| {sep} |")
            continue

        # key:value 쌍 묶기
        pairs = []
        i = 0
        while i < len(cells):
            col, val = cells[i]
            if i + 1 < len(cells):
                next_col, next_val = cells[i + 1]
                if next_col - col <= 2:
                    pairs.append(f"**{val}**: {next_val}")
                    i += 2
                    continue
            pairs.append(val)
            i += 1
        parts.append("- " + " / ".join(pairs))

    return "\n".join(parts)


def _split_column_regions(non_empty_cols: list[int],
                          use_cols: list[int] | None = None,
                          raw_rows: list[list[str]] | None = None,
                          gap_threshold: int = 3,
                          min_region_cols: int = 3) -> list[list[int]]:
    """비어있지 않은 열 인덱스(raw_rows 기준)를 원본 워크시트 열 갭 기준으로 분리한다.

    Args:
        non_empty_cols: raw_rows 내 데이터가 있는 열 인덱스 (0-based)
        use_cols: raw_rows 열 → 원본 워크시트 열 매핑 (1-based)
        raw_rows: 원본 행 데이터 (소규모 영역 병합 판단용)
        gap_threshold: 원본 워크시트 열 번호 기준 이 값 이상 갭이면 분리
        min_region_cols: 이보다 적은 열의 영역은 인접 영역에 병합

    Returns:
        [[raw_col_indices], ...] — 분리된 영역별 raw 열 인덱스 리스트
    """
    if not non_empty_cols:
        return []

    regions = [[non_empty_cols[0]]]
    for i in range(1, len(non_empty_cols)):
        cur_raw = non_empty_cols[i]
        prev_raw = non_empty_cols[i - 1]
        if use_cols:
            ws_gap = use_cols[cur_raw] - use_cols[prev_raw]
        else:
            ws_gap = cur_raw - prev_raw
        if ws_gap >= gap_threshold:
            regions.append([])
        regions[-1].append(cur_raw)

    # 영역이 1개면 분리 불필요
    if len(regions) <= 1:
        return regions

    # 소규모 영역 병합: 데이터가 거의 없는 영역만 다음 영역에 합침
    if raw_rows:
        merged = []
        pending_small = []
        for reg in regions:
            # 영역 내 비어있지 않은 행 수 계산
            nonempty_rows = sum(
                1 for row in raw_rows
                if any(row[c] for c in reg)
            )
            if nonempty_rows < 3 and len(reg) < min_region_cols:
                pending_small.append(reg)
            else:
                if pending_small:
                    for small in pending_small:
                        reg = small + reg
                    pending_small = []
                merged.append(reg)
        if pending_small:
            if merged:
                for small in pending_small:
                    merged[-1].extend(small)
            else:
                merged = [non_empty_cols]
        return merged

    return regions


def _deduplicate_headers(headers: list[str]) -> tuple[list[str], list[int]]:
    """중복 헤더 텍스트를 처리하고 유니크한 열 이름을 반환한다.

    task-127 Step 6 시정: 중복 헤더 열을 단순 제거하는 대신 suffix(_2, _3, ...)를
    붙여서 데이터를 보존한다. 예: ['변화율', '변화율'] → ['변화율', '변화율_2'].
    이렇게 하면 기간별 변화율 컬럼처럼 의미있는 중복 컬럼의 데이터 손실을 방지한다.

    수평 병합(h_merges) 기반 중복만 제거 — 이는 호출부에서 use_cols 레벨에서 이미
    처리되므로 여기서는 suffix 방식으로 전환한다.

    Returns:
        (renamed_headers, keep_col_indices)
        keep_col_indices: 변경이 있으면 전체 인덱스 목록, 없으면 [] (호환성 유지)
    """
    seen: dict[str, int] = {}  # 헤더 값 → 마지막 suffix 번호
    renamed: list[str] = []
    has_change = False

    for i, h in enumerate(headers):
        if h and h in seen:
            # 중복 헤더 → suffix 증가하여 보존
            seen[h] += 1
            renamed.append(f"{h}_{seen[h]}")
            has_change = True
        else:
            seen[h] = 1
            renamed.append(h)

    if not has_change:
        return headers, []  # 중복 없음 — 원본 반환, keep_cols=[] (전체 열 유지)

    # 중복이 있으면 keep_cols = 전체 인덱스 (모든 열 보존)
    return renamed, list(range(len(renamed)))


def _format_pre_header(pre_header: list[list[str]]) -> str:
    """헤더 전 메타데이터 행을 요약 텍스트로 변환한다."""
    lines = []
    for row in pre_header:
        non_empty = [v for v in row if v]
        if not non_empty:
            continue
        if len(non_empty) == 1:
            lines.append(f"**{non_empty[0]}**")
        elif len(non_empty) == 2:
            lines.append(f"- {non_empty[0]}: {non_empty[1]}")
        elif len(non_empty) <= 4:
            lines.append("- " + " / ".join(non_empty))
        else:
            # task-127 S4-3: 5+ 셀 보존 + MAX_PRE_HEADER_CELLS_PER_ROW=50 가드
            shown = non_empty[:MAX_PRE_HEADER_CELLS_PER_ROW]
            remainder = len(non_empty) - len(shown)
            line = "- " + " / ".join(shown)
            if remainder > 0:
                line += f" _( 외 {remainder}개 셀 생략)_"
            lines.append(line)
    return "\n".join(lines)


def _detect_header_row(rows: list[list[str]]) -> Optional[int]:
    """테이블 헤더 행을 감지한다.

    우선순위:
    1. 'ID', '번호', 'No' 등 테이블 헤더 키워드가 포함된 행
    2. 3개 이상 고유 값이 있고, 모두 숫자가 아닌 첫 행
    """
    HEADER_KEYWORDS = {'id', '번호', 'no', 'no.', '항목', '분류', 'category',
                       'name', 'type', 'status', '결과', '내용', 'component',
                       '구분', 'description', 'title', 'device', 'os'}

    # 1차: 키워드 매칭
    for i, row in enumerate(rows):
        non_empty = [v for v in row if v]
        if len(non_empty) >= 3:
            lower_vals = {v.lower().strip() for v in non_empty}
            if lower_vals & HEADER_KEYWORDS:
                return i

    # 2차: 3개 이상 비숫자 고유값
    for i, row in enumerate(rows):
        non_empty = [v for v in row if v]
        if len(non_empty) >= 3:
            all_numeric = all(
                re.match(r'^-?[\d,.]+%?$', v) for v in non_empty
            )
            if not all_numeric:
                unique_ratio = len(set(non_empty)) / len(non_empty)
                if unique_ratio >= 0.5:  # 절반 이상 고유값
                    return i

    return None


def _fill_down_categories(headers: list[str], rows: list[list[str]]) -> list[list[str]]:
    """분류/카테고리 성격의 열에서 빈 셀을 위 행 값으로 채운다.

    판단 기준: 열 이름에 '분류', 'category', '구분', 'type', 'Component' 등이 포함되거나,
    해당 열의 빈 셀 비율이 50% 이상이고 첫 데이터가 비어있지 않은 경우.
    """
    CATEGORY_KEYWORDS = {'분류', 'category', '구분', 'type', 'component',
                         '대분류', '중분류', '소분류', 'group', '항목'}

    category_cols = set()
    for c_idx, header in enumerate(headers):
        h_lower = header.lower().strip()
        if any(kw in h_lower for kw in CATEGORY_KEYWORDS):
            category_cols.add(c_idx)

    # 빈 셀 비율 기반 감지 (50% 이상 빈 셀 + 첫 행에 값 있음)
    if rows:
        for c_idx in range(len(headers)):
            if c_idx in category_cols:
                continue
            empty_count = sum(1 for r in rows if not r[c_idx])
            total = len(rows)
            if total > 5 and empty_count / total > 0.5:
                # 첫 값이 있고, 값이 반복적 패턴인 경우
                first_val = rows[0][c_idx] if rows[0][c_idx] else None
                if first_val:
                    category_cols.add(c_idx)

    if not category_cols:
        return rows

    # fill-down 실행
    result = []
    prev_vals = [""] * len(headers)
    for row in rows:
        new_row = list(row)
        for c_idx in category_cols:
            if c_idx < len(new_row):
                if new_row[c_idx]:
                    prev_vals[c_idx] = new_row[c_idx]
                else:
                    new_row[c_idx] = prev_vals[c_idx]
        result.append(new_row)

    return result


def _sanitize_cells(cells: list[str]) -> list[str]:
    """셀 값을 마크다운 테이블에 안전하게 넣을 수 있도록 정리한다.

    task-127 v3 시정 2: 절단 시 끝부분 issue ID 패턴(#XXXXXX)이 있으면 보존.
    multi-line 셀 ('... \\n#162733')이 MAX_COL_WIDTH 초과로 끝 ID 손실되던 문제 시정.
    """
    result = []
    for v in cells:
        v = v.replace("|", "\\|").replace("\n", " ").replace("\r", "")
        # MINOR-4 (task-125): 수식 augment 셀은 절단 임계값 동적 확장
        # 형태: "<eval> [=...]" → 끝이 "]"이고 " [=" 패턴 포함
        is_formula_augmented = v.endswith("]") and " [=" in v
        max_w = MAX_COL_WIDTH if not is_formula_augmented else max(MAX_COL_WIDTH, len(v))
        if len(v) > max_w:
            # task-127 v3 시정 2: 끝부분에 ID 패턴이 있으면 보존
            tail_ids = _TRAILING_ID_RE.findall(v)
            if tail_ids:
                # set으로 중복 제거 후 끝에 부착 (검색 토큰 노출 우선)
                tail_str = " " + " ".join(dict.fromkeys(tail_ids))  # 순서 보존 dedup
                head_w = max(0, max_w - len(tail_str))
                v = v[:head_w] + "…" + tail_str
            else:
                v = v[:max_w] + "…"
        if not v:
            v = " "  # 빈 셀은 공백으로 (테이블 깨짐 방지)
        result.append(v)
    return result


# ═══════════════════════════════════════════════════════════════════════════════
# XLSX 이미지/수식 helper (task-125 R-1, R-2)
# ═══════════════════════════════════════════════════════════════════════════════

def _extract_xlsx_image(
    image_obj,
    sheet_name: str,
    img_seq: int,
    img_dir: Optional[Path],
    save: bool,
) -> Optional[dict]:
    """xlsx 이미지에서 anchor + bytes를 추출하고 선택적으로 디스크에 저장.

    Args:
        image_obj: openpyxl Image 객체
        sheet_name: 시트명 (marker A1 표기용)
        img_seq: 전체 이미지 시퀀스 번호 (1-based)
        img_dir: 저장 디렉터리 (None이면 메모리만)
        save: 디스크 저장 여부

    Returns:
        {filename, content_type, size_bytes, sheet, anchor_cell, anchor_row,
         anchor_col, saved_path, format} 또는 실패 시 None
    """
    try:
        from openpyxl.utils import get_column_letter

        # MINOR-5 (task-125): AbsoluteAnchor fallback — hasattr 분기 필수
        anchor = image_obj.anchor
        if hasattr(anchor, '_from') and anchor._from is not None:
            # TwoCellAnchor / OneCellAnchor
            row_0 = anchor._from.row
            col_0 = anchor._from.col
            anchor_row = row_0 + 1     # 1-based
            anchor_col = col_0 + 1     # 1-based
            anchor_cell = f"{get_column_letter(anchor_col)}{anchor_row}"
        else:
            # AbsoluteAnchor (pos.x/pos.y는 EMU 좌표 — 셀 변환 불가)
            anchor_row = None
            anchor_col = None
            anchor_cell = "?"
            logger.debug(
                "[file_parsers] xlsx anchor type=%s — anchor_cell='?' fallback",
                type(anchor).__name__,
            )

        # blob 추출
        data_attr = getattr(image_obj, '_data', None)
        if callable(data_attr):
            blob = data_attr()
        else:
            blob = data_attr or b""
        if not blob:
            logger.warning(
                "[file_parsers] xlsx image blob 비어있음 sheet=%s seq=%d",
                sheet_name, img_seq,
            )
            return None

        # format + filename
        ext = (getattr(image_obj, 'format', None) or 'png').lower()
        safe_sheet = _safe_filename(sheet_name)
        filename = f"{safe_sheet}_img{img_seq:03d}.{ext}"
        content_type = f"image/{'jpeg' if ext == 'jpg' else ext}"

        # save (옵션)
        saved_path = None
        if save and img_dir:
            try:
                img_dir.mkdir(parents=True, exist_ok=True)
                save_path = img_dir / filename
                save_path.write_bytes(blob)
                saved_path = str(save_path)
            except Exception as e:
                logger.warning(
                    "[file_parsers] xlsx image 저장 실패 sheet=%s: %s",
                    sheet_name, e,
                )

        return {
            "filename": filename,
            "content_type": content_type,
            "size_bytes": len(blob),
            "sheet": sheet_name,
            "anchor_cell": anchor_cell,
            "anchor_row": anchor_row,
            "anchor_col": anchor_col,
            "saved_path": saved_path,
            "format": ext,
        }
    except Exception as e:
        logger.warning(
            "[file_parsers] xlsx image 추출 실패 sheet=%s seq=%d: %s",
            sheet_name, img_seq, e,
        )
        return None


def _extract_xlsx_image_section(
    ws,
    sheet_name: str,
    img_dir: Optional[Path],
    save: bool,
    images_accumulator: list,
    total_images_offset: int,
) -> tuple:
    """시트의 모든 이미지를 추출하고 marker 라인 + OCR 텍스트를 반환.

    C-3 시정 (task-125): summary/non-summary 두 경로 모두 호출되도록 분리.

    Args:
        ws: openpyxl worksheet
        sheet_name: 시트명 (marker A1 표기용)
        img_dir: 이미지 저장 디렉터리 (None이면 메모리만)
        save: 디스크 저장 여부
        images_accumulator: parse_xlsx의 images 리스트 (mutate)
        total_images_offset: 누적 이미지 수 (img_seq 시작값)

    Returns:
        (marker_lines, count) — marker_lines는 body_text에 추가할 문자열 리스트,
        count는 처리한 이미지 수.
    """
    if not hasattr(ws, "_images") or not ws._images:
        return [], 0

    marker_lines: list[str] = []
    count = 0
    for local_seq, image_obj in enumerate(ws._images, 1):
        img_seq = total_images_offset + local_seq
        info = _extract_xlsx_image(image_obj, sheet_name, img_seq, img_dir, save)
        if info:
            images_accumulator.append(info)
            count += 1
            size_kb = info["size_bytes"] // 1024
            marker = (
                f'[이미지: {sheet_name}!{info["anchor_cell"]}, '
                f'{info["filename"]}, {size_kb}KB]'
            )
            marker_lines.append(marker)
            # OCR 옵션 — 저장 경로가 있는 경우만
            if info.get("saved_path"):
                ocr_text = _ocr_image(info["saved_path"])
                if ocr_text:
                    marker_lines.append(ocr_text)
    return marker_lines, count


def _load_xlsx_formula_map(file_path: str) -> dict:
    """수식 셀 좌표 → 원문 매핑을 read_only 모드로 빠르게 추출.

    R-2 / 옵션 D (task-125): 메인 data_only=True wb는 그대로 두고,
    보조 read_only=False wb로 수식 원문만 수집. 처리시간 추가 ~1초 (47MB 기준).

    Returns:
        {sheet_name: {(row_1based, col_1based): "=SUM(A1:A10)", ...}}
        실패 시 빈 dict (graceful degradation — 평가값만 보존).
    """
    import openpyxl
    formula_map: dict = {}
    try:
        wb_r = openpyxl.load_workbook(file_path, data_only=False, read_only=True)
        for sname in wb_r.sheetnames:
            ws_r = wb_r[sname]
            sheet_map: dict = {}
            for row in ws_r.iter_rows():
                for cell in row:
                    # data_type='f' = formula cell
                    if cell.data_type == 'f' and isinstance(cell.value, str):
                        sheet_map[(cell.row, cell.column)] = cell.value
            if sheet_map:
                formula_map[sname] = sheet_map
        wb_r.close()
    except Exception as e:
        logger.warning(
            "[file_parsers] xlsx formula map 추출 실패: %s — 평가값만 보존", e,
        )
    return formula_map


def _format_cell_value_with_formula(
    val: str,
    formula: Optional[str],
) -> str:
    """평가값에 수식 원문을 부착 (있으면).

    R-2 / M-4 (task-125): step2 포맷 채택 — 'value [=FORMULA]'
    평가값 우선, 수식 괄호 부착. 검색 토큰 노출 + 가독성 균형.

    Args:
        val: 평가값 (data_only=True 결과, 이미 str+strip)
        formula: formula_map.get((r, c)) 결과 (None이면 일반 셀)

    Returns:
        formula 있으면 'val [=FORMULA]', 없으면 val 그대로.
    """
    if not formula:
        return val
    if not val:
        # 평가값 없는 수식 셀 (희소) — 수식만 표시
        return formula
    # 검색 토큰 노출 + 가독성 균형 — 평가값 우선, 수식은 괄호 부착
    return f"{val} [{formula}]"


# ═══════════════════════════════════════════════════════════════════════════════
# PDF 파서 (task-125 R-3)
# ═══════════════════════════════════════════════════════════════════════════════

def parse_pdf(file_path: str) -> dict:
    """PDF 파일을 페이지별 마크다운 텍스트 + 표로 변환한다.

    1차 범위:
      - 페이지별 텍스트 추출 (## Page N 헤더)
      - 표 추출 → 마크다운 테이블 (pdfplumber.extract_tables)
    제외 (후속 task escalate):
      - OCR (스캔 PDF) — is_scanned_likely=True 로그만
      - 이미지 추출
      - 수식 (LaTeX)

    Returns:
        {
            "body_text": str,
            "metadata": {
                "file_path": str,
                "file_name": str,
                "total_pages": int,
                "total_tables": int,
                "is_scanned_likely": bool,
            },
            "images": [],
            "source_type": "generic_pdf",
        }
    """
    try:
        import pdfplumber
    except ImportError:
        logger.warning(
            "[file_parsers] pdfplumber 미설치 — PDF 파싱 불가 (%s)", file_path,
        )
        return {
            "body_text": "[pdfplumber 미설치 — PDF 파싱 불가]",
            "metadata": {"file_path": file_path},
            "images": [],
            "source_type": "generic_pdf",
        }

    file_name = Path(file_path).name
    parts = []
    total_pages = 0
    total_tables = 0
    text_total_chars = 0

    try:
        with pdfplumber.open(file_path) as pdf:
            total_pages = len(pdf.pages)
            for page_idx, page in enumerate(pdf.pages, 1):
                page_parts = [f"## Page {page_idx}"]
                try:
                    page_text = page.extract_text() or ""
                except Exception as e:
                    logger.debug("page %d text extract 실패: %s", page_idx, e)
                    page_text = ""
                page_text = page_text.strip()
                if page_text:
                    page_parts.append(page_text)
                    text_total_chars += len(page_text)
                try:
                    tables = page.extract_tables() or []
                except Exception as e:
                    logger.debug("page %d tables extract 실패: %s", page_idx, e)
                    tables = []
                for table in tables:
                    md = _pdf_table_to_markdown(table)
                    if md:
                        page_parts.append(md)
                        total_tables += 1
                if len(page_parts) > 1:
                    parts.append("\n".join(page_parts))
    except Exception as e:
        logger.error("[file_parsers] pdfplumber open 실패 (%s): %s", file_path, e)
        return {
            "body_text": "",
            "metadata": {
                "file_path": file_path,
                "file_name": file_name,
                "error": f"{type(e).__name__}: {e}",
            },
            "images": [],
            "source_type": "generic_pdf",
        }

    body_text = "\n\n".join(parts)
    if len(body_text) > MAX_BODY_CHARS:
        body_text = body_text[:MAX_BODY_CHARS] + "\n\n_(본문 잘림)_"

    is_scanned_likely = (total_pages > 0 and text_total_chars < total_pages * 10)
    if is_scanned_likely and total_pages > 0:
        logger.warning(
            "[file_parsers] PDF likely scanned (no extractable text): "
            "%s (pages=%d, total_chars=%d)",
            file_path, total_pages, text_total_chars,
        )

    return {
        "body_text": body_text,
        "metadata": {
            "file_path": file_path,
            "file_name": file_name,
            "total_pages": total_pages,
            "total_tables": total_tables,
            "is_scanned_likely": is_scanned_likely,
        },
        "images": [],
        "source_type": "generic_pdf",
    }


def _pdf_table_to_markdown(table: list) -> str:
    """pdfplumber 테이블 (list of rows) → 마크다운 테이블.

    Args:
        table: list[list[Optional[str]]] — pdfplumber extract_tables() 결과

    Returns:
        마크다운 테이블 문자열. 빈 표이면 '' 반환.
    """
    if not table or not table[0]:
        return ""
    sanitized = []
    for row in table:
        cells = []
        for v in row:
            v = "" if v is None else str(v).strip()
            v = v.replace("|", "\\|").replace("\n", " ")
            if len(v) > MAX_COL_WIDTH:
                v = v[:MAX_COL_WIDTH] + "…"
            if not v:
                v = " "
            cells.append(v)
        sanitized.append(cells)

    # 모두 빈 경우
    if not any(any(c.strip() for c in r) for r in sanitized):
        return ""

    n_cols = len(sanitized[0])
    out = ["| " + " | ".join(sanitized[0]) + " |",
           "|" + "|".join("---" for _ in sanitized[0]) + "|"]
    for row in sanitized[1:]:
        # 열 수 정규화
        while len(row) < n_cols:
            row.append(" ")
        row = row[:n_cols]
        out.append("| " + " | ".join(row) + " |")
    return "\n".join(out)


# ═══════════════════════════════════════════════════════════════════════════════
# TSV 파서
# ═══════════════════════════════════════════════════════════════════════════════

def parse_tsv(file_path: str) -> dict:
    """TSV 파일을 마크다운 테이블로 변환한다.

    헤더(첫 행) + 데이터 행을 마크다운 테이블로 출력.
    큰 파일은 MAX_TABLE_ROWS로 잘라낸다.
    """
    file_name = Path(file_path).name
    rows = []
    headers = []
    encoding = _detect_encoding(file_path)

    try:
        # 대용량 TSV 필드 지원 (스토리/대사 등 긴 텍스트)
        csv.field_size_limit(10 * 1024 * 1024)  # 10MB
        with open(file_path, "r", encoding=encoding, errors="replace") as f:
            # BOM 제거
            reader = csv.reader(f, delimiter="\t")
            for i, row in enumerate(reader):
                # 빈 행 건너뛰기
                if not any(cell.strip() for cell in row):
                    continue
                if i == 0:
                    headers = [h.strip() for h in row]
                else:
                    rows.append([cell.strip() for cell in row])
    except Exception as e:
        logger.error("TSV 파일 읽기 실패: %s — %s", file_path, e)
        return {
            "body_text": f"[TSV 파일 읽기 실패: {e}]",
            "metadata": {"file_path": file_path, "error": str(e)},
            "images": [],
            "source_type": "generic_tsv",
        }

    if not headers:
        return {
            "body_text": "[TSV 파일 비어있음]",
            "metadata": {"file_path": file_path},
            "images": [],
            "source_type": "generic_tsv",
        }

    # 빈 열 제거
    non_empty_cols = []
    for c in range(len(headers)):
        if headers[c] or any(c < len(r) and r[c] for r in rows):
            non_empty_cols.append(c)

    if non_empty_cols:
        headers = [headers[c] for c in non_empty_cols]
        rows = [[r[c] if c < len(r) else "" for c in non_empty_cols] for r in rows]

    # 마크다운 테이블 생성
    parts = []
    sanitized_h = _sanitize_cells(headers)
    parts.append("| " + " | ".join(sanitized_h) + " |")
    parts.append("|" + "|".join("---" for _ in sanitized_h) + "|")

    truncated = False
    for i, row in enumerate(rows):
        if i >= MAX_TABLE_ROWS:
            truncated = True
            break
        # 열 수 맞추기
        while len(row) < len(headers):
            row.append("")
        parts.append("| " + " | ".join(_sanitize_cells(row[:len(headers)])) + " |")

    if truncated:
        parts.append(f"\n_(… 외 {len(rows) - MAX_TABLE_ROWS}행 생략, 총 {len(rows)}행)_")

    body_text = "\n".join(parts)
    if len(body_text) > MAX_BODY_CHARS:
        body_text = body_text[:MAX_BODY_CHARS] + "\n\n_(본문 잘림)_"

    return {
        "body_text": body_text,
        "metadata": {
            "file_path": file_path,
            "file_name": file_name,
            "total_rows": len(rows),
            "total_columns": len(headers),
            "headers": headers,
        },
        "images": [],
        "source_type": "generic_tsv",
    }


def _detect_encoding(file_path: str) -> str:
    """파일 인코딩을 감지한다 (BOM 기반 + fallback)."""
    with open(file_path, "rb") as f:
        raw = f.read(4)
    if raw[:3] == b"\xef\xbb\xbf":
        return "utf-8-sig"
    if raw[:2] in (b"\xff\xfe", b"\xfe\xff"):
        return "utf-16"
    return "utf-8"


# ═══════════════════════════════════════════════════════════════════════════════
# PNG/이미지 파서
# ═══════════════════════════════════════════════════════════════════════════════

def parse_image(file_path: str) -> dict:
    """이미지 파일의 메타데이터를 추출한다.

    Vision API 연동은 별도 호출 (여기서는 메타데이터만).
    """
    file_name = Path(file_path).name
    file_size = os.path.getsize(file_path)

    width, height = 0, 0
    img_format = Path(file_path).suffix.lstrip(".")

    try:
        from PIL import Image
        with Image.open(file_path) as img:
            width, height = img.size
            img_format = img.format or img_format
    except ImportError:
        logger.info("Pillow 미설치 — 이미지 크기 추출 불가")
    except Exception as e:
        logger.warning("이미지 메타 추출 실패: %s — %s", file_path, e)

    size_kb = file_size // 1024
    body_text = (
        f"## 이미지: {file_name}\n"
        f"- 형식: {img_format.upper()}\n"
        f"- 크기: {width}x{height} px\n"
        f"- 파일 크기: {size_kb} KB\n"
        f"- 경로: {file_path}\n"
    )

    return {
        "body_text": body_text,
        "metadata": {
            "file_path": file_path,
            "file_name": file_name,
            "width": width,
            "height": height,
            "format": img_format,
            "file_size_bytes": file_size,
        },
        "images": [{"filename": file_name, "saved_path": file_path,
                     "size_bytes": file_size, "content_type": f"image/{img_format.lower()}"}],
        "source_type": "image",
    }


# ═══════════════════════════════════════════════════════════════════════════════
# DOCX 파서
# ═══════════════════════════════════════════════════════════════════════════════

def parse_docx(file_path: str) -> dict:
    """DOCX 파일을 마크다운 텍스트로 변환한다.

    헤딩 스타일은 마크다운 헤딩(##)으로, 테이블은 마크다운 테이블로 변환.
    이미지는 참조 텍스트만 남긴다.
    """
    try:
        from docx import Document
    except ImportError:
        return {
            "body_text": "[python-docx 미설치 — DOCX 파싱 불가]",
            "metadata": {"file_path": file_path},
            "images": [],
            "source_type": "generic_docx",
        }

    doc = Document(file_path)
    parts: list[str] = []
    file_name = Path(file_path).name

    # 헤딩 스타일 → 마크다운 레벨 매핑
    heading_map = {
        "Heading 1": "#", "Heading 2": "##", "Heading 3": "###",
        "Heading 4": "####", "Heading 5": "#####", "Heading 6": "######",
        "Title": "#", "Subtitle": "##",
    }

    for element in doc.element.body:
        tag = element.tag.split("}")[-1]  # namespace 제거

        if tag == "p":
            # 단락 처리
            from docx.text.paragraph import Paragraph
            para = Paragraph(element, doc)
            text = para.text.strip()
            if not text:
                continue

            style_name = para.style.name if para.style else ""
            prefix = heading_map.get(style_name, "")
            if prefix:
                parts.append(f"{prefix} {text}")
            else:
                # 불릿/번호 리스트 감지
                numPr = element.find(
                    ".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}numPr"
                )
                if numPr is not None:
                    parts.append(f"- {text}")
                else:
                    parts.append(text)

        elif tag == "tbl":
            # 테이블 처리
            from docx.table import Table
            table = Table(element, doc)
            rows_data = []
            for row in table.rows:
                cells = [cell.text.replace("|", "\\|").replace("\n", " ").strip()
                         for cell in row.cells]
                rows_data.append(cells)

            if rows_data:
                # 첫 행을 헤더로
                header = rows_data[0]
                parts.append("| " + " | ".join(h or " " for h in header) + " |")
                parts.append("|" + "|".join("---" for _ in header) + "|")
                for row_cells in rows_data[1:]:
                    # 열 수 맞추기
                    while len(row_cells) < len(header):
                        row_cells.append(" ")
                    parts.append(
                        "| " + " | ".join(c or " " for c in row_cells[:len(header)]) + " |"
                    )

    body_text = "\n\n".join(parts)
    # NBSP 등 특수 공백 정리
    body_text = body_text.replace("\xa0", " ").replace("\u200b", "")
    if len(body_text) > MAX_BODY_CHARS:
        body_text = body_text[:MAX_BODY_CHARS] + "\n\n_(본문 잘림)_"

    return {
        "body_text": body_text,
        "metadata": {
            "file_path": file_path,
            "file_name": file_name,
            "total_paragraphs": len(doc.paragraphs),
            "total_tables": len(doc.tables),
        },
        "images": [],
        "source_type": "generic_docx",
    }


# ═══════════════════════════════════════════════════════════════════════════════
# HTML 파서
# ═══════════════════════════════════════════════════════════════════════════════

def parse_html(file_path: str, *, extract_images: bool = True,
               image_base_dir: Optional[str] = None) -> dict:
    """HTML 파일을 마크다운 텍스트로 변환한다.

    맥락 보존 우선순위:
      1) 헤딩 구조 (h1~h6) → 마크다운 헤딩
      2) 문단/리스트/표 → 마크다운
      3) <img alt/title> 텍스트 (캡션 역할) 보존
      4) <figure>/<figcaption> 보존
      5) base64 인라인 이미지 OCR (Tesseract 가용 시)
      6) <script>/<style>/<nav>/<footer> 제거

    외부 src URL 이미지는 다운로드하지 않음 (보안/속도). alt/title만 유지.
    """
    try:
        from bs4 import BeautifulSoup, NavigableString
    except ImportError:
        return {
            "body_text": "[BeautifulSoup 미설치 — HTML 파싱 불가]",
            "metadata": {"file_path": file_path},
            "images": [],
            "source_type": "generic_html",
        }

    import base64
    import binascii

    # 인코딩 자동 감지 (UTF-8 우선, 실패 시 CP949)
    raw_bytes = Path(file_path).read_bytes()
    for enc in ("utf-8", "cp949", "euc-kr", "latin-1"):
        try:
            html = raw_bytes.decode(enc)
            break
        except UnicodeDecodeError:
            continue
    else:
        html = raw_bytes.decode("utf-8", errors="replace")

    soup = BeautifulSoup(html, "html.parser")

    # 1) 불필요 태그 제거
    for tag_name in ("script", "style", "nav", "footer", "noscript", "template"):
        for t in soup.find_all(tag_name):
            t.decompose()

    file_name = Path(file_path).name
    parts: list[str] = []
    images_info: list[dict] = []
    total_images = 0
    ocr_count = 0

    # 이미지 저장 디렉터리 준비 (base64 추출용)
    if extract_images:
        if image_base_dir:
            img_dir = Path(image_base_dir) / IMAGE_DIR_NAME / _safe_filename(Path(file_path).stem)
        else:
            img_dir = Path(file_path).parent / IMAGE_DIR_NAME / _safe_filename(Path(file_path).stem)
    else:
        img_dir = None

    # 제목 보존
    title_tag = soup.find("title")
    if title_tag and title_tag.get_text(strip=True):
        parts.append(f"# {title_tag.get_text(strip=True)}")

    body = soup.body or soup

    def _flatten_text(node) -> str:
        """블록 내 인라인 텍스트를 공백 연결하여 반환."""
        return " ".join(node.get_text(separator=" ", strip=True).split())

    def _render_table(table_tag) -> str:
        rows: list[list[str]] = []
        for tr in table_tag.find_all("tr"):
            cells = [
                c.get_text(separator=" ", strip=True).replace("|", "\\|")
                for c in tr.find_all(["th", "td"])
            ]
            if cells:
                rows.append(cells)
        if not rows:
            return ""
        # 첫 행 헤더 가정
        out = ["| " + " | ".join(rows[0]) + " |",
               "|" + "|".join("---" for _ in rows[0]) + "|"]
        for row in rows[1:]:
            # 열 수 맞춤 (rowspan/colspan 정확 처리는 생략 — 채우기)
            pad = len(rows[0]) - len(row)
            if pad > 0:
                row = row + [""] * pad
            elif pad < 0:
                row = row[:len(rows[0])]
            out.append("| " + " | ".join(row) + " |")
        return "\n".join(out)

    def _handle_image(img_tag) -> Optional[str]:
        """img 태그에서 맥락 정보 추출. alt/title 우선, base64는 OCR 시도."""
        nonlocal total_images, ocr_count
        total_images += 1
        alt = (img_tag.get("alt") or "").strip()
        title = (img_tag.get("title") or "").strip()
        src = (img_tag.get("src") or "").strip()

        ref_bits = []
        if alt:
            ref_bits.append(f"alt={alt}")
        if title and title != alt:
            ref_bits.append(f"title={title}")

        # base64 인라인 이미지 처리
        ocr_text = ""
        if src.startswith("data:image/") and "base64," in src and img_dir:
            try:
                header, b64 = src.split("base64,", 1)
                mime = header[5:].split(";")[0]  # e.g. "image/png"
                ext = mime.split("/")[-1] or "png"
                img_bytes = base64.b64decode(b64)
                img_dir.mkdir(parents=True, exist_ok=True)
                fname = f"html_img{total_images:03d}.{ext}"
                save_path = img_dir / fname
                save_path.write_bytes(img_bytes)
                images_info.append({
                    "filename": fname,
                    "content_type": mime,
                    "size_bytes": len(img_bytes),
                    "saved_path": str(save_path),
                })
                # OCR
                if extract_images:
                    ocr_text = _ocr_image(str(save_path))
                    if ocr_text:
                        ocr_count += 1
            except (binascii.Error, ValueError) as e:
                logger.debug("base64 이미지 디코드 실패: %s", e)

        if not ref_bits and not ocr_text:
            return None

        label = ", ".join(ref_bits) if ref_bits else src[:40] + "…"
        out = f"[이미지: {label}]"
        if ocr_text:
            out += f"\n> OCR: {ocr_text.strip()}"
        return out

    def _walk(node, depth: int = 0):
        """DOM을 순회하며 parts에 마크다운을 추가한다."""
        if isinstance(node, NavigableString):
            text = str(node).strip()
            if text:
                parts.append(text)
            return

        if node.name is None:
            return

        name = node.name.lower()

        if name in ("h1", "h2", "h3", "h4", "h5", "h6"):
            level = int(name[1])
            text = _flatten_text(node)
            if text:
                parts.append(f"{'#' * level} {text}")
            return

        if name == "p":
            text = _flatten_text(node)
            if text:
                parts.append(text)
            return

        if name in ("ul", "ol"):
            for li in node.find_all("li", recursive=False):
                t = _flatten_text(li)
                if t:
                    bullet = "-" if name == "ul" else "1."
                    parts.append(f"{bullet} {t}")
            return

        if name == "table":
            md = _render_table(node)
            if md:
                parts.append(md)
            return

        if name == "img":
            out = _handle_image(node)
            if out:
                parts.append(out)
            return

        if name == "figure":
            # figure 안의 img 먼저 처리 (figcaption은 별도로)
            for child in node.children:
                if hasattr(child, "name") and child.name == "figcaption":
                    continue
                _walk(child, depth + 1)
            caption = node.find("figcaption")
            if caption:
                cap_text = _flatten_text(caption)
                if cap_text:
                    parts.append(f"*Caption: {cap_text}*")
            return

        if name == "pre":
            # 코드 블록
            code_text = node.get_text()
            if code_text.strip():
                parts.append("```\n" + code_text.rstrip() + "\n```")
            return

        if name == "blockquote":
            text = _flatten_text(node)
            if text:
                parts.append(f"> {text}")
            return

        if name in ("br", "hr"):
            parts.append("")
            return

        # 기본: 자식 노드 재귀
        for child in node.children:
            _walk(child, depth + 1)

    _walk(body)

    body_text = "\n\n".join(p for p in parts if p.strip())
    if len(body_text) > MAX_BODY_CHARS:
        body_text = body_text[:MAX_BODY_CHARS] + "\n\n_(본문 잘림)_"

    return {
        "body_text": body_text,
        "metadata": {
            "file_path": file_path,
            "file_name": file_name,
            "total_images": total_images,
            "ocr_hits": ocr_count,
            "title": title_tag.get_text(strip=True) if title_tag else "",
        },
        "images": images_info,
        "source_type": "generic_html",
    }


# ═══════════════════════════════════════════════════════════════════════════════
# 유틸리티
# ═══════════════════════════════════════════════════════════════════════════════

def _safe_filename(name: str) -> str:
    """파일명에서 위험한 문자를 제거한다."""
    return re.sub(r'[<>:"/\\|?*]', "_", name)


# ── CLI 테스트 ────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("Usage: python file_parsers.py <file_path>")
        sys.exit(1)

    path = sys.argv[1]
    result = parse_file(path, extract_images=False)

    print(f"\n{'='*60}")
    print(f"source_type: {result['source_type']}")
    print(f"metadata: {result['metadata']}")
    print(f"images: {len(result['images'])}개")
    print(f"body_text 길이: {len(result['body_text'])} chars")
    print(f"{'='*60}")
    print(result["body_text"][:3000])
    if len(result["body_text"]) > 3000:
        print(f"\n... (총 {len(result['body_text'])} chars)")
